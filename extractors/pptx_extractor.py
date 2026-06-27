from pptx import Presentation

def extract_pptx(file_path):
    prs = Presentation(file_path)
    pages = []
    for slide_num, slide in enumerate(prs.slides, 1):
        lines = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = para.text.strip()
                    if text:
                        lines.append(text)
        pages.append((slide_num, "\n".join(lines)))
    return pages
